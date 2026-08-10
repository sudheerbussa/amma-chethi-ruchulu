#!/usr/bin/env python3
"""Generic WhatsApp Order & Booking System pitch — Storyset-style illustrated version.

Keeps the plain deck untouched:
  pitch/WhatsApp_Order_System_Client_Pitch.pptx
Writes:
  pitch/WhatsApp_Order_System_Client_Pitch_Storyset.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "pitch" / "WhatsApp_Order_System_Client_Pitch_Storyset.pptx"
ASSETS = Path(__file__).resolve().parents[1] / "pitch" / "storyset-assets"
OFFICIAL = ASSETS / "official"


def img(name: str) -> Path:
    """Prefer official Storyset downloads; fall back to local assets."""
    o = OFFICIAL / name
    if o.exists():
        return o
    return ASSETS / name

NAVY = RGBColor(0x1A, 0x2B, 0x3C)
TEAL = RGBColor(0x0D, 0x7A, 0x6F)
TEAL_LIGHT = RGBColor(0xE6, 0xF4, 0xF2)
WARM = RGBColor(0xC4, 0x5C, 0x26)
OFF_WHITE = RGBColor(0xF7, 0xF5, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5A, 0x64, 0x6E)
DARK = RGBColor(0x1F, 0x29, 0x37)


def set_run(run, text, size=18, bold=False, color=DARK, font="Calibri", hyperlink=None):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    if hyperlink:
        run.hyperlink.address = hyperlink


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    font="Calibri",
    hyperlink=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        # Only attach hyperlink to the first line unless text is a single URL line
        link = hyperlink if (hyperlink and (i == 0 or "wa.me" in line.lower() or line.startswith("http"))) else None
        if hyperlink and len(lines) == 1:
            link = hyperlink
        set_run(p.add_run(), line, size, bold, color, font, hyperlink=link)
    return box


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_shape(shape, color)
    return shape


def add_round_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shape, color)
    return shape


def footer(slide, page, total):
    add_textbox(
        slide,
        Inches(0.5),
        Inches(7.15),
        Inches(9),
        Inches(0.3),
        "WhatsApp Order & Booking System  ·  Confidential",
        size=10,
        color=GRAY,
    )
    add_textbox(
        slide,
        Inches(11.2),
        Inches(7.15),
        Inches(1.5),
        Inches(0.3),
        f"{page} / {total}",
        size=10,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )


def title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.15), NAVY)
    add_textbox(
        slide,
        Inches(0.55),
        Inches(0.28),
        Inches(12),
        Inches(0.45),
        title,
        size=24,
        bold=True,
        color=WHITE,
        font="Georgia",
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.55),
            Inches(0.72),
            Inches(12),
            Inches(0.35),
            subtitle,
            size=12,
            color=RGBColor(0xB8, 0xC5, 0xD0),
        )


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def table_row(slide, x0, y, col_w, cells, bg, color=DARK, bold=False, size=12, height=0.55):
    for i, cell in enumerate(cells):
        x = x0 + sum(col_w[:i])
        add_rect(slide, Inches(x), Inches(y), Inches(col_w[i]), Inches(height), bg)
        add_textbox(
            slide,
            Inches(x + 0.08),
            Inches(y + 0.12),
            Inches(col_w[i] - 0.12),
            Inches(height - 0.1),
            cell,
            size=size,
            bold=bold,
            color=color,
            align=PP_ALIGN.CENTER if i else PP_ALIGN.LEFT,
        )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 15

    # --- 1 Title ---
    s = blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), TEAL)
    add_textbox(
        s,
        Inches(0.7),
        Inches(0.55),
        Inches(7),
        Inches(0.4),
        "WhatsApp Order & Booking System",
        size=14,
        bold=True,
        color=TEAL,
    )
    add_textbox(
        s,
        Inches(0.7),
        Inches(1.8),
        Inches(7),
        Inches(1.4),
        "Orders, appointments\n& payments on WhatsApp",
        size=32,
        bold=True,
        color=WHITE,
        font="Georgia",
    )
    add_textbox(
        s,
        Inches(0.7),
        Inches(3.5),
        Inches(6.8),
        Inches(1.0),
        "Built for India’s small local businesses.\nNo app download for your customers.",
        size=16,
        color=RGBColor(0xC9, 0xD4, 0xDE),
    )
    add_textbox(
        s,
        Inches(0.7),
        Inches(4.8),
        Inches(6.8),
        Inches(0.4),
        "Food · shops · clinics · home services",
        size=14,
        color=TEAL,
    )
    add_textbox(s, Inches(0.7), Inches(6.5), Inches(7), Inches(0.35), "Illustrated pitch  ·  2026", size=12, color=GRAY)
    if img("01-hero-mobile-marketing.png").exists():
        s.shapes.add_picture(str(img("01-hero-mobile-marketing.png")), Inches(7.4), Inches(1.5), width=Inches(5.4))

    # --- 2 Problem ---
    s = blank_slide(prs)
    title_bar(s, "The problem", "Phone calls and messy WhatsApp chats don’t scale")
    if img("02-problem-questions.png").exists():
        s.shapes.add_picture(str(img("02-problem-questions.png")), Inches(0.5), Inches(1.5), height=Inches(5.2))
    problems = [
        ("Missed messages", "Orders and booking requests get buried."),
        ("No structure", "No menu, slots, payment trail, or status."),
        ("Aggregator tax", "Food brands lose 18–30% on apps."),
        ("Hard to build alone", "Official WhatsApp API is not DIY."),
    ]
    for i, (h, body) in enumerate(problems):
        top = Inches(1.5 + i * 1.25)
        add_round_rect(s, Inches(6.6), top, Inches(6.2), Inches(1.1), OFF_WHITE)
        add_textbox(s, Inches(6.85), top + Inches(0.18), Inches(5.7), Inches(0.35), h, size=16, bold=True, color=TEAL, font="Georgia")
        add_textbox(s, Inches(6.85), top + Inches(0.55), Inches(5.7), Inches(0.4), body, size=13, color=DARK)
    footer(s, 2, total)

    # --- 3 Solution ---
    s = blank_slide(prs)
    title_bar(s, "The solution", "Your WhatsApp number becomes the storefront / booking desk")
    add_textbox(
        s,
        Inches(0.6),
        Inches(1.35),
        Inches(12),
        Inches(0.4),
        "Customer chats → picks items or a time slot → pays (UPI/Razorpay) → you get a clear ops board.",
        size=15,
        color=DARK,
    )
    if img("03-solution-chatbot.png").exists():
        s.shapes.add_picture(str(img("03-solution-chatbot.png")), Inches(3.5), Inches(1.9), height=Inches(4.0))
    add_textbox(
        s,
        Inches(0.6),
        Inches(6.15),
        Inches(12),
        Inches(0.6),
        "Official Meta WhatsApp Cloud API only. Two product lines: Order (menu + cart) and Booking (appointments / slots).",
        size=13,
        color=GRAY,
    )
    footer(s, 3, total)

    # --- 4 Who ---
    s = blank_slide(prs)
    title_bar(s, "Who this is for", "Small local businesses already living on WhatsApp")
    # collage of 3 official storyset images (left column)
    for i, name in enumerate(["04-verticals-coffee.png", "04b-doctor.png", "04c-groceries.png"]):
        if img(name).exists():
            s.shapes.add_picture(str(img(name)), Inches(0.45), Inches(1.4 + i * 1.8), height=Inches(1.7))
    cols = [
        ("Order-first", "Restaurants, tiffin, bakery, meat, flower, kirana (phase 2)"),
        ("Booking-first", "Clinics, dental, salons, tuition slots, auto service, studios"),
        ("Repeat-service", "Laundry, water can, AC service, pest control, cleaning, tailor"),
    ]
    for i, (h, body) in enumerate(cols):
        top = Inches(1.5 + i * 1.7)
        add_round_rect(s, Inches(6.5), top, Inches(6.3), Inches(1.5), OFF_WHITE)
        add_rect(s, Inches(6.5), top, Inches(0.12), Inches(1.5), [NAVY, TEAL, WARM][i])
        add_textbox(s, Inches(6.85), top + Inches(0.25), Inches(5.7), Inches(0.4), h, size=18, bold=True, color=NAVY, font="Georgia")
        add_textbox(s, Inches(6.85), top + Inches(0.7), Inches(5.7), Inches(0.6), body, size=13, color=DARK)
    footer(s, 4, total)

    # --- 5 What you get ---
    s = blank_slide(prs)
    title_bar(s, "What you get", "Not just a chatbot greeting — a working ops system")
    features = [
        ("WhatsApp storefront", "Buttons, lists, cart / slots, confirmations"),
        ("Payments", "UPI QR + Razorpay checkout (where needed)"),
        ("Ops dashboard", "Orders / bookings board, status, CSV export"),
        ("Customer memory", "Saved address / last booking, re-order / re-book"),
        ("Alerts", "Staff notifications + optional kitchen/desk voice alerts"),
        ("Hosting included", "Shared VPS + SSL on All-in Hosted plans"),
    ]
    for i, (h, b) in enumerate(features):
        col, row = i % 3, i // 3
        left, top = Inches(0.5 + col * 4.2), Inches(1.55 + row * 2.55)
        add_round_rect(s, left, top, Inches(3.95), Inches(2.25), OFF_WHITE)
        add_rect(s, left, top, Inches(0.12), Inches(2.25), TEAL)
        add_textbox(s, left + Inches(0.35), top + Inches(0.4), Inches(3.4), Inches(0.45), h, size=17, bold=True, color=NAVY, font="Georgia")
        add_textbox(s, left + Inches(0.35), top + Inches(1.0), Inches(3.4), Inches(0.9), b, size=14, color=GRAY)
    footer(s, 5, total)

    # --- 6 Plans ---
    s = blank_slide(prs)
    title_bar(s, "Plans for small businesses (India)", "All-in Hosted · GST extra · Meta marketing messages billed separately")
    if img("05-plans-business.png").exists():
        s.shapes.add_picture(str(img("05-plans-business.png")), Inches(9.7), Inches(1.5), width=Inches(3.2))
    headers = ["", "Starter", "Growth", "Business"]
    rows = [
        ["Setup (one-time)", "₹4,999", "₹7,999", "₹12,999"],
        ["Monthly (hosted)", "₹1,499", "₹2,499", "₹3,499"],
        ["Per paid order / booking", "₹8", "₹6", "₹5"],
        ["Catalog / menu items*", "Up to 40", "Up to 100", "Up to 250"],
        ["Go-live", "7–10 days", "10–14 days", "2–3 weeks"],
    ]
    col_w = [2.8, 2.1, 2.1, 2.1]
    x0 = 0.4
    table_row(s, x0, 1.4, col_w, headers, NAVY, WHITE, True, 11, 0.45)
    for r, row in enumerate(rows):
        bg = TEAL_LIGHT if r % 2 == 0 else OFF_WHITE
        for i, cell in enumerate(row):
            x = x0 + sum(col_w[:i])
            cell_bg = NAVY if i == 0 else bg
            add_rect(s, Inches(x), Inches(1.85 + r * 0.62), Inches(col_w[i]), Inches(0.62), cell_bg)
            add_textbox(
                s,
                Inches(x + 0.08),
                Inches(1.97 + r * 0.62),
                Inches(col_w[i] - 0.12),
                Inches(0.4),
                cell,
                size=11,
                bold=(i == 0),
                color=WHITE if i == 0 else DARK,
                align=PP_ALIGN.CENTER if i else PP_ALIGN.LEFT,
            )
    add_round_rect(s, Inches(0.4), Inches(5.2), Inches(8.9), Inches(1.4), TEAL_LIGHT)
    add_textbox(
        s,
        Inches(0.6),
        Inches(5.4),
        Inches(8.5),
        Inches(1.05),
        "Custom plans for bigger businesses (retail 250+ SKUs, multi-outlet, multi-doctor clinics, ERP) — contact us.\n"
        "*Booking: “items” = services / doctors / rooms. Hosted = shared VPS + SSL. Meta marketing & Maps add-ons extra.",
        size=12,
        color=DARK,
    )
    footer(s, 6, total)

    # --- 7 Why plans cost more (KEY SLIDE) ---
    s = blank_slide(prs)
    title_bar(s, "Why plans cost more as you go up", "Price follows complexity — not a random jump")
    metrics = [
        ("Catalog size", "More items = more categories, stock rules, testing, and WhatsApp list design."),
        ("Order / booking volume", "Higher traffic needs stronger hosting share, monitoring, and support response."),
        ("Payment & offers", "UPI-only is simple; Razorpay + offers + reminders add build and ops work."),
        ("Staff & locations", "1 phone vs multi-staff / multi-counter needs roles, alerts, and training."),
        ("Custom rules", "Cutoffs, advance pay, doctor-wise slots, delivery zones — each rule is extra build."),
        ("Support depth", "Starter = email/WhatsApp help; Business includes faster response + more changes."),
    ]
    for i, (h, b) in enumerate(metrics):
        col, row = i % 3, i // 3
        left, top = Inches(0.45 + col * 4.2), Inches(1.4 + row * 2.6)
        add_round_rect(s, left, top, Inches(4.0), Inches(2.35), OFF_WHITE)
        add_textbox(s, left + Inches(0.25), top + Inches(0.3), Inches(3.5), Inches(0.4), h, size=16, bold=True, color=TEAL, font="Georgia")
        add_textbox(s, left + Inches(0.25), top + Inches(0.85), Inches(3.5), Inches(1.2), b, size=13, color=DARK)
    footer(s, 7, total)

    # --- 8 Plan limits comparison ---
    s = blank_slide(prs)
    title_bar(s, "Plan limits at a glance", "So you can match your shop size to a plan")
    headers2 = ["Limit", "Starter", "Growth", "Business"]
    rows2 = [
        ["Menu / service items", "40", "100", "250"],
        ["Categories", "5", "10", "20"],
        ["Staff WhatsApp admins", "1", "3", "6"],
        ["Locations / counters", "1", "1", "Up to 3"],
        ["Razorpay checkout", "Add-on / UPI", "Included", "Included"],
        ["Offers / broadcasts*", "Basic", "Yes", "Yes + priority"],
        ["Training sessions", "1", "2", "3"],
        ["Hypercare after go-live", "7 days", "14 days", "21 days"],
    ]
    col_w2 = [3.5, 2.8, 2.8, 2.8]
    table_row(s, 0.55, 1.4, col_w2, headers2, NAVY, WHITE, True, 12, 0.45)
    for r, row in enumerate(rows2):
        y = 1.85 + r * 0.55
        for i, cell in enumerate(row):
            x = 0.55 + sum(col_w2[:i])
            bg = NAVY if i == 0 else (TEAL_LIGHT if r % 2 == 0 else OFF_WHITE)
            add_rect(s, Inches(x), Inches(y), Inches(col_w2[i]), Inches(0.55), bg)
            add_textbox(
                s,
                Inches(x + 0.08),
                Inches(y + 0.12),
                Inches(col_w2[i] - 0.1),
                Inches(0.35),
                cell,
                size=12,
                bold=(i == 0),
                color=WHITE if i == 0 else DARK,
                align=PP_ALIGN.CENTER if i else PP_ALIGN.LEFT,
            )
    add_textbox(
        s,
        Inches(0.55),
        Inches(6.4),
        Inches(12),
        Inches(0.4),
        "*Broadcasts use Meta marketing rates (client-paid). Item limits keep WhatsApp UX usable and setup effort predictable.",
        size=11,
        color=GRAY,
    )
    footer(s, 8, total)

    # --- 9 What's included vs extra ---
    s = blank_slide(prs)
    title_bar(s, "What’s included vs extra cost", "No surprises after you say yes")
    left_items = [
        "Shared VPS hosting + SSL",
        "WhatsApp Cloud API setup help",
        "Menu / services load (within plan limit)",
        "UPI payment path",
        "Admin dashboard",
        "Basic training + go-live support",
    ]
    right_items = [
        "Meta message fees (esp. marketing)",
        "Own custom domain (optional)",
        "One-page website if you don’t have one",
        "Google Maps auto-distance",
        "Items beyond plan limit",
        "Custom workflows / integrations",
    ]
    add_round_rect(s, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.1), TEAL_LIGHT)
    add_textbox(s, Inches(0.8), Inches(1.75), Inches(5.3), Inches(0.4), "Included in All-in Hosted", size=18, bold=True, color=NAVY, font="Georgia")
    for i, t in enumerate(left_items):
        add_textbox(s, Inches(0.9), Inches(2.4 + i * 0.55), Inches(5.2), Inches(0.45), "▸  " + t, size=14, color=DARK)
    add_round_rect(s, Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.1), OFF_WHITE)
    add_textbox(s, Inches(7.2), Inches(1.75), Inches(5.3), Inches(0.4), "Extra / client-paid", size=18, bold=True, color=WARM, font="Georgia")
    for i, t in enumerate(right_items):
        add_textbox(s, Inches(7.3), Inches(2.4 + i * 0.55), Inches(5.2), Inches(0.45), "▸  " + t, size=14, color=DARK)
    footer(s, 9, total)

    # --- 10 Clinics & booking ---
    s = blank_slide(prs)
    title_bar(s, "Booking businesses (clinics & more)", "Same WhatsApp channel — slots instead of a food menu")
    add_textbox(
        s,
        Inches(0.6),
        Inches(1.35),
        Inches(7),
        Inches(0.45),
        "Patient / client picks doctor or service → date/time slot → optional advance pay → reminder.",
        size=14,
        color=DARK,
    )
    booking = [
        ("Clinics", "GP, dental, physio — cut front-desk phone load"),
        ("Diagnostics", "Lab booking + report-ready update"),
        ("Salons / spas", "Stylist + slot + festival advance"),
        ("Coaching", "Demo class / batch seat booking"),
        ("Workshops", "Limited seats, waitlist, reminders"),
        ("Service bays", "Car service appointment windows"),
    ]
    for i, (h, b) in enumerate(booking):
        col, row = i % 2, i // 2
        left, top = Inches(0.5 + col * 3.7), Inches(2.0 + row * 1.45)
        add_round_rect(s, left, top, Inches(3.5), Inches(1.3), OFF_WHITE)
        add_textbox(s, left + Inches(0.2), top + Inches(0.2), Inches(3.1), Inches(0.35), h, size=14, bold=True, color=TEAL, font="Georgia")
        add_textbox(s, left + Inches(0.2), top + Inches(0.6), Inches(3.1), Inches(0.55), b, size=12, color=DARK)
    if img("09-online-doctor.png").exists():
        s.shapes.add_picture(str(img("09-online-doctor.png")), Inches(8.2), Inches(1.9), width=Inches(4.6))
    footer(s, 10, total)

    # --- 11 Case study ---
    s = blank_slide(prs)
    title_bar(s, "Case study: Amma Chethi Ruchulu", "Live reference — food ordering on official WhatsApp")
    add_round_rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.0), TEAL_LIGHT)
    add_textbox(
        s,
        Inches(0.9),
        Inches(1.75),
        Inches(11.5),
        Inches(0.4),
        "Tenali home-cooked Andhra food brand — production WhatsApp ordering",
        size=18,
        bold=True,
        color=NAVY,
        font="Georgia",
    )
    add_textbox(
        s,
        Inches(0.9),
        Inches(2.35),
        Inches(11.5),
        Inches(0.8),
        "Lunch/dinner cutoffs, multi-cook menu, Razorpay/UPI, kitchen dashboard, delivery status — "
        "running on Meta Cloud API (order.ammachethiruchulu.co.in).",
        size=14,
        color=DARK,
    )
    bullets = [
        "Proof that the stack works in a real kitchen — not a slideware demo",
        "Your business gets a white-label version matched to your menu or booking rules",
        "Live demo shown on the discovery call (we do not publish the ACR order number in the deck)",
    ]
    for i, b in enumerate(bullets):
        add_textbox(s, Inches(0.9), Inches(3.9 + i * 0.65), Inches(11.5), Inches(0.5), "▸  " + b, size=15, color=DARK)
    footer(s, 11, total)

    # --- 12 Engagement ---
    s = blank_slide(prs)
    title_bar(s, "How we engage", "From first call to first paid order / booking")
    steps2 = [
        ("01", "Discovery", "Menu or services, volumes, staff, website status"),
        ("02", "Plan match", "Starter / Growth / Business / Custom using your limits"),
        ("03", "WhatsApp setup", "WABA, number, Meta verification support"),
        ("04", "Configure & train", "Load catalog/slots, dry-run with staff"),
        ("05", "Go live", "Hypercare window based on plan"),
    ]
    for i, (n, t, b) in enumerate(steps2):
        top = Inches(1.45 + i * 0.95)
        add_round_rect(s, Inches(0.55), top, Inches(12.2), Inches(0.85), TEAL_LIGHT if i % 2 == 0 else OFF_WHITE)
        add_textbox(s, Inches(0.8), top + Inches(0.22), Inches(0.8), Inches(0.45), n, size=20, bold=True, color=TEAL, font="Georgia")
        add_textbox(s, Inches(1.8), top + Inches(0.15), Inches(3.5), Inches(0.55), t, size=16, bold=True, color=NAVY)
        add_textbox(s, Inches(5.6), top + Inches(0.22), Inches(6.8), Inches(0.5), b, size=14, color=GRAY)
    footer(s, 12, total)

    # --- 13 What we need from you ---
    s = blank_slide(prs)
    title_bar(
        s,
        "What we need from you",
        "Same checklist we use for official WhatsApp + online payments (India)",
    )
    cols_need = [
        (
            "Business registration",
            [
                "Udyam certificate (PDF) — MSME",
                "PAN (proprietor / business)",
                "GSTIN if you have it",
                "Shop & Establishment (optional)",
                "Legal name + address (exact match)",
            ],
        ),
        (
            "WhatsApp / Meta",
            [
                "Dedicated +91 mobile number",
                "Number free for Cloud API OTP",
                "Business email + phone",
                "Website or we build one-pager",
                "Privacy policy URL (required)",
            ],
        ),
        (
            "Payments (Razorpay / PhonePe)",
            [
                "Business / current bank account",
                "Cancelled cheque or bank proof",
                "Aadhaar (prop.) if gateway asks",
                "Business proof (Udyam / GST)",
                "UPI VPA if starting UPI-only",
            ],
        ),
    ]
    for i, (h, items) in enumerate(cols_need):
        left = Inches(0.4 + i * 4.25)
        add_round_rect(s, left, Inches(1.4), Inches(4.1), Inches(4.55), OFF_WHITE)
        add_rect(s, left, Inches(1.4), Inches(4.1), Inches(0.6), [NAVY, TEAL, WARM][i])
        add_textbox(
            s,
            left + Inches(0.15),
            Inches(1.52),
            Inches(3.8),
            Inches(0.4),
            h,
            size=14,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        for j, item in enumerate(items):
            add_textbox(
                s,
                left + Inches(0.25),
                Inches(2.2 + j * 0.65),
                Inches(3.7),
                Inches(0.55),
                "•  " + item,
                size=13,
                color=DARK,
            )
    add_textbox(
        s,
        Inches(0.5),
        Inches(6.15),
        Inches(12.3),
        Inches(0.6),
        "Names/address must match across Udyam, Meta, website, and payment gateway. We guide uploads; approval is by Meta / Razorpay / PhonePe.\n"
        "Also share: menu or service list, delivery/booking rules, staff WhatsApp numbers for admin alerts.",
        size=12,
        color=GRAY,
    )
    footer(s, 13, total)

    # --- 14 Add-ons ---
    s = blank_slide(prs)
    title_bar(s, "Optional add-ons", "Buy only what you need")
    addons = [
        ("Meta-ready one-page website", "₹2,999 – ₹4,999", "If you don’t have a site yet"),
        ("Own domain connect", "₹999 + domain cost", "Brand URL instead of subdomain"),
        ("Google Maps distance", "₹999/mo or ₹0.50–1/order", "Auto km fee instead of zone buttons"),
        ("Extra items pack (+50)", "₹1,499 one-time", "When you outgrow plan catalog"),
        ("Extra training session", "₹999", "New staff / festival rush"),
        ("Custom workflow", "₹1,500 – ₹2,500 / hour", "Special rules beyond template"),
    ]
    for i, (h, price, note) in enumerate(addons):
        col, row = i % 3, i // 3
        left, top = Inches(0.45 + col * 4.2), Inches(1.5 + row * 2.55)
        add_round_rect(s, left, top, Inches(4.0), Inches(2.3), OFF_WHITE)
        add_textbox(s, left + Inches(0.25), top + Inches(0.3), Inches(3.5), Inches(0.7), h, size=15, bold=True, color=NAVY, font="Georgia")
        add_textbox(s, left + Inches(0.25), top + Inches(1.1), Inches(3.5), Inches(0.4), price, size=16, bold=True, color=TEAL)
        add_textbox(s, left + Inches(0.25), top + Inches(1.55), Inches(3.5), Inches(0.5), note, size=12, color=GRAY)
    footer(s, 14, total)

    # --- 15 CTA ---
    s = blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), TEAL)
    add_textbox(
        s,
        Inches(0.7),
        Inches(1.2),
        Inches(7.5),
        Inches(0.9),
        "Ready to run orders or bookings on WhatsApp?",
        size=26,
        bold=True,
        color=WHITE,
        font="Georgia",
    )
    add_textbox(
        s,
        Inches(0.7),
        Inches(2.2),
        Inches(7.5),
        Inches(0.8),
        "Share your menu size or appointment volume — we’ll map you to\nStarter, Growth, Business, or a Custom plan.",
        size=15,
        color=RGBColor(0xC9, 0xD4, 0xDE),
    )
    WA_URL = (
        "https://wa.me/918055292935?text="
        "Hi%2C%20I%20want%20to%20book%20a%20discovery%20call%20for%20WhatsApp%20Order%20System"
    )
    btn = add_round_rect(s, Inches(0.7), Inches(3.3), Inches(5.6), Inches(1.15), TEAL)
    btn.click_action.hyperlink.address = WA_URL
    tf = btn.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    set_run(p0.add_run(), "Book a discovery call", size=14, bold=True, color=WHITE)
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    set_run(p1.add_run(), "+91 80552 92935", size=22, bold=True, color=WHITE)
    qr = ASSETS / "wa-discovery-qr.png"
    if qr.exists():
        add_round_rect(s, Inches(0.7), Inches(5.15), Inches(1.55), Inches(1.55), WHITE)
        s.shapes.add_picture(str(qr), Inches(0.8), Inches(5.25), width=Inches(1.35), height=Inches(1.35))
        add_textbox(
            s,
            Inches(2.4),
            Inches(5.55),
            Inches(5.5),
            Inches(0.9),
            "Scan to WhatsApp\nDemo on discovery call · ACR case study",
            size=12,
            bold=True,
            color=RGBColor(0xC9, 0xD4, 0xDE),
        )
    if img("06-cta-contact.png").exists():
        s.shapes.add_picture(str(img("06-cta-contact.png")), Inches(8.2), Inches(1.4), width=Inches(4.6))
    add_textbox(
        s,
        Inches(0.7),
        Inches(6.85),
        Inches(12),
        Inches(0.3),
        "Illustrations: Storyset by Freepik (storyset.com) — free use with attribution",
        size=10,
        color=GRAY,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({total} slides)")


if __name__ == "__main__":
    build()
